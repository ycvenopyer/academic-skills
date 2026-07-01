#!/usr/bin/env python3
"""
论文汇报PPT生成工具
基于提取的论文JSON数据，使用 python-pptx 直接生成 PPTX 文件。
适合组会论文汇报场景。

用法:
  python paper_ppt.py paper_ppt_data.json [--output presentation.pptx] [--template academic]
"""

import argparse
import json
import os
import re
import sys
from pathlib import Path


def create_presentation(slides_data: dict, output_path: str,
                        template_style: str = "academic",
                        figures_dir: str = None,
                        tables_dir: str = None):
    """
    基于幻灯片数据创建PPTX文件
    模板: academic（蓝色学术风格，默认）
    figures_dir: extract-picture 提取的图片目录（自动插入对应图表）
    tables_dir: extract-picture 提取的表格目录
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    import copy

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    # === 颜色方案 ===
    COLORS = {
        "academic": {
            "primary": RGBColor(0x1B, 0x3A, 0x5C),     # 深蓝
            "secondary": RGBColor(0x2C, 0x5F, 0x8A),    # 中蓝
            "accent": RGBColor(0x3B, 0x82, 0xC4),       # 亮蓝
            "light_bg": RGBColor(0xE8, 0xF0, 0xF8),     # 浅蓝背景
            "white": RGBColor(0xFF, 0xFF, 0xFF),
            "dark_text": RGBColor(0x2D, 0x2D, 0x2D),
            "gray_text": RGBColor(0x66, 0x66, 0x66),
            "light_gray": RGBColor(0xF5, 0xF5, 0xF5),
            "accent_gold": RGBColor(0xC9, 0x9A, 0x3B),  # 金色强调
        }
    }

    C = COLORS.get(template_style, COLORS["academic"])
    paper_info = slides_data.get("paper_info", {})
    slides = slides_data.get("slides", [])

    def add_text_box(slide, left, top, width, height, text, font_size=18,
                     bold=False, color=C["dark_text"], alignment=PP_ALIGN.LEFT,
                     font_name="微软雅黑"):
        """添加文本框"""
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = alignment
        return txBox

    def add_bilingual_content(slide, left, top, width, height, original, translation,
                              orig_size=12, trans_size=16, font_name="微软雅黑"):
        """添加双语对照内容（左侧原文，右侧译文）"""
        half_w = width / 2 - 0.2
        # 原文（左侧，灰色小字）
        txBox_orig = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(half_w), Inches(height))
        tf_orig = txBox_orig.text_frame
        tf_orig.word_wrap = True
        p = tf_orig.paragraphs[0]
        p.text = f"[原文]\n{original[:800]}"
        p.font.size = Pt(orig_size)
        p.font.color.rgb = C["gray_text"]
        p.font.name = font_name

        # 译文（右侧，深色大字）
        txBox_trans = slide.shapes.add_textbox(
            Inches(left + half_w + 0.3), Inches(top), Inches(half_w), Inches(height))
        tf_trans = txBox_trans.text_frame
        tf_trans.word_wrap = True
        p = tf_trans.paragraphs[0]
        p.text = f"[译文]\n{translation[:800]}"
        p.font.size = Pt(trans_size)
        p.font.color.rgb = C["dark_text"]
        p.font.name = font_name
        p.font.bold = True

    def find_figure_for_slide(slide_title: str, figures_dir: str) -> str:
        """根据幻灯片标题在 figures_dir 中查找匹配的图片"""
        if not figures_dir or not os.path.isdir(figures_dir):
            return None

        # 从标题提取关键词
        title_lower = slide_title.lower()
        keywords = []
        for kw in ["method", "model", "architecture", "experiment", "result",
                    "comparison", "overview", "framework", "pipeline",
                    "方法", "模型", "架构", "实验", "结果", "对比", "框架"]:
            if kw in title_lower:
                keywords.append(kw)

        # 读取 manifest 找匹配的图片
        manifest_path = os.path.join(figures_dir, "_manifest.md")
        image_candidates = []
        if os.path.exists(manifest_path):
            with open(manifest_path, 'r', encoding='utf-8') as f:
                for line in f:
                    line = line.strip()
                    # 匹配表格行中的图片文件名
                    if '|' in line and ('.png' in line or '.jpg' in line):
                        cells = [c.strip() for c in line.split('|') if c.strip()]
                        for cell in cells:
                            if cell.endswith(('.png', '.jpg', '.jpeg')):
                                fpath = os.path.join(figures_dir, cell)
                                if os.path.exists(fpath):
                                    image_candidates.append((fpath, cell, line))

        # 有关键词时，优先找标题中匹配的
        if keywords and image_candidates:
            for fpath, fname, manifest_line in image_candidates:
                manifest_lower = manifest_line.lower()
                if any(kw in manifest_lower for kw in keywords):
                    return fpath

        # 无关键词匹配或没有 manifest 时，按排序取第一张非图标图片
        for f in sorted(os.listdir(figures_dir)):
            if f.endswith(('.png', '.jpg', '.jpeg')) and not f.startswith('_'):
                # 跳过可能的小图标（文件名含 icon/logo）
                if any(tag in f.lower() for tag in ['icon', 'logo', 'arrow']):
                    continue
                return os.path.join(figures_dir, f)

        # fallback: 取任意第一张
        if image_candidates:
            return image_candidates[0][0]

        return None

    def add_image_to_slide(slide, img_path, left, top, width, height):
        """在幻灯片中插入图片，自动缩放保持比例"""
        try:
            from pptx.util import Inches
            slide.shapes.add_picture(img_path, Inches(left), Inches(top),
                                     Inches(width), Inches(height))
            return True
        except Exception as e:
            print(f"  插入图片失败 {img_path}: {e}")
            return False

    def add_multiline_text(slide, left, top, width, height, lines, font_size=16,
                           color=C["dark_text"], line_spacing=1.5, font_name="微软雅黑"):
        """添加多行文本（支持列表）"""
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.name = font_name
            p.space_after = Pt(font_size * 0.5)
        return txBox

    def add_bg_rect(slide, left, top, width, height, color):
        """添加背景色块"""
        from pptx.util import Inches
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()  # 无边框
        return shape

    # === 生成幻灯片 ===
    for slide_idx, slide_data in enumerate(slides):
        slide_type = slide_data.get("type", "content")
        title = slide_data.get("title", "")

        # 添加空白幻灯片
        slide_layout = prs.slide_layouts[6]  # blank layout
        slide = prs.slides.add_slide(slide_layout)

        if slide_type == "title":
            # 标题页 - 深蓝背景
            add_bg_rect(slide, 0, 0, 13.333, 7.5, C["primary"])
            # 装饰条
            add_bg_rect(slide, 0, 3.0, 13.333, 0.08, C["accent_gold"])

            add_text_box(slide, 1.5, 1.5, 10.3, 2.0, title,
                         font_size=36, bold=True, color=C["white"],
                         alignment=PP_ALIGN.CENTER)
            add_text_box(slide, 1.5, 3.5, 10.3, 1.0,
                         slide_data.get("subtitle", ""),
                         font_size=18, color=C["accent"],
                         alignment=PP_ALIGN.CENTER)

        elif slide_type == "outline":
            # 目录页
            add_bg_rect(slide, 0, 0, 13.333, 1.5, C["primary"])
            add_text_box(slide, 1.0, 0.3, 11.3, 1.0, title,
                         font_size=32, bold=True, color=C["white"])
            add_bg_rect(slide, 0, 1.5, 13.333, 0.05, C["accent_gold"])

            items = slide_data.get("items", [])
            add_multiline_text(slide, 2.0, 2.0, 9.3, 5.0,
                               [f"  {i+1}. {item}" for i, item in enumerate(items)],
                               font_size=22, color=C["dark_text"])

        elif slide_type == "content":
            # 内容页
            add_bg_rect(slide, 0, 0, 13.333, 1.2, C["primary"])
            add_text_box(slide, 0.8, 0.2, 11.7, 0.8, title,
                         font_size=28, bold=True, color=C["white"])
            add_bg_rect(slide, 0, 1.2, 13.333, 0.05, C["accent_gold"])

            is_bilingual = slide_data.get("bilingual", False)

            if is_bilingual:
                # 双语布局：左侧原文，右侧译文
                original = slide_data.get("original", slide_data.get("content", ""))
                translation = slide_data.get("translation", "")
                add_bilingual_content(slide, 0.8, 1.6, 11.7, 5.0,
                                      original[:1500], translation[:1500])
            else:
                content = slide_data.get("content", "")
                subslides = slide_data.get("subslides", [])
                sub_idx = slide_data.get("sub_index", 0)

                display_content = content
                if subslides and len(subslides) > sub_idx:
                    display_content = subslides[sub_idx]

                # 处理段落
                paragraphs = [p.strip() for p in display_content.split('\n') if p.strip()]
                short_paras = []
                for p in paragraphs[:15]:
                    if len(p) > 200:
                        sentences = re.split(r'(?<=[.?!])\s+', p)
                        short_paras.extend(sentences[:3])
                    else:
                        short_paras.append(p)

                # 判断是否有图片可插入
                img_path = find_figure_for_slide(title, figures_dir) if figures_dir else None
                if img_path and len(short_paras) < 6:
                    # 左文右图布局
                    add_multiline_text(slide, 0.8, 1.6, 6.5, 5.5, short_paras[:8],
                                       font_size=16, color=C["dark_text"])
                    add_image_to_slide(slide, img_path, 7.8, 1.8, 4.5, 4.5)
                else:
                    # 纯文字布局
                    add_multiline_text(slide, 0.8, 1.6, 11.7, 5.5, short_paras[:12],
                                       font_size=16, color=C["dark_text"])

            # 页码
            if slide_idx < len(slides):
                add_text_box(slide, 12.0, 7.0, 1.0, 0.4,
                             f"{slide_idx+1}/{len(slides)}",
                             font_size=10, color=C["gray_text"],
                             alignment=PP_ALIGN.RIGHT)

        elif slide_type == "ending":
            # 总结页
            add_bg_rect(slide, 0, 0, 13.333, 1.5, C["primary"])
            add_text_box(slide, 1.0, 0.3, 11.3, 1.0, title,
                         font_size=32, bold=True, color=C["white"])
            add_bg_rect(slide, 0, 1.5, 13.333, 0.05, C["accent_gold"])

            points = slide_data.get("points", [])
            add_multiline_text(slide, 1.5, 2.0, 10.3, 5.0,
                               [f"💡 {p}" for p in points],
                               font_size=20, color=C["dark_text"])

        elif slide_type == "qa":
            # Q&A页
            add_bg_rect(slide, 0, 0, 13.333, 7.5, C["primary"])
            add_bg_rect(slide, 5.5, 3.2, 2.3, 0.08, C["accent_gold"])
            add_text_box(slide, 1.5, 2.5, 10.3, 2.0, "谢谢！",
                         font_size=48, bold=True, color=C["white"],
                         alignment=PP_ALIGN.CENTER)
            add_text_box(slide, 1.5, 4.0, 10.3, 1.0, "欢迎提问 🙋",
                         font_size=24, color=C["accent"],
                         alignment=PP_ALIGN.CENTER)

    prs.save(output_path)
    return True


def main():
    parser = argparse.ArgumentParser(description="论文汇报PPT生成工具")
    parser.add_argument("input", help="输入JSON文件（来自 extract_paper_for_ppt.py）")
    parser.add_argument("--output", "-o", default=None, help="输出PPTX文件路径")
    parser.add_argument("--template", choices=["academic"], default="academic",
                        help="PPT模板风格")
    parser.add_argument("--figures-dir", default=None,
                        help="extract-picture 提取的图片目录，自动插入对应幻灯片")
    parser.add_argument("--tables-dir", default=None,
                        help="extract-picture 提取的表格目录（用于表格页）")

    args = parser.parse_args()

    if not os.path.exists(args.input):
        print(f"错误: 文件不存在 - {args.input}")
        sys.exit(1)

    with open(args.input, 'r', encoding='utf-8') as f:
        slides_data = json.load(f)

    if args.output is None:
        input_stem = Path(args.input).stem.replace("_ppt_data", "").replace("_data", "")
        args.output = f"{input_stem}_汇报PPT.pptx"

    print(f"正在生成PPT: {args.output}")
    print(f"   幻灯片数: {len(slides_data.get('slides', []))}")
    print(f"   论文: {slides_data.get('paper_info', {}).get('title', '未知')[:60]}")

    if args.figures_dir:
        if os.path.isdir(args.figures_dir):
            figure_files = [f for f in os.listdir(args.figures_dir)
                            if f.endswith(('.png', '.jpg', '.jpeg')) and not f.startswith('_')]
            print(f"   图片目录: {args.figures_dir} ({len(figure_files)}张图片)")
        else:
            print(f"   警告: 图片目录不存在 - {args.figures_dir}")

    create_presentation(slides_data, args.output, args.template,
                        figures_dir=args.figures_dir,
                        tables_dir=args.tables_dir)

    print(f"\n✅ PPT生成完成: {os.path.abspath(args.output)}")
    print(f"\n提示：")
    print(f"   1. 如需要调整内容，编辑JSON文件后重新运行")
    print(f"   2. 插入图片: 使用 --figures-dir 指向 extract-picture 的输出目录")
    print(f"   3. 双语PPT: 在JSON中将 slide.bilingual 设为 true 并填写 translation 字段")
    print(f"   4. 如需更精美的PPT，使用 html2pptx 工作流（见 SKILL.md）")
    print(f"   5. 生成的PPT可在PowerPoint/WPS中进一步微调")


if __name__ == "__main__":
    main()
